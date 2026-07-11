#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from run_v5 results.

Uses /usr/bin/python3 (Python 3.9) which has python-pptx installed.
Run with:
    /usr/bin/python3 scripts/generate_ppt.py
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

REPO_ROOT  = Path(__file__).resolve().parents[1]
RUN_V5     = REPO_ROOT / "artifacts" / "run_v5"
COMBINED   = RUN_V5 / "combined_report"
OUT_PATH   = RUN_V5 / "results.pptx"

RUNS = [
    {"dataset": "adult",   "noise_pct": "20%", "dir": RUN_V5 / "adult_20pct"},
    {"dataset": "german",  "noise_pct": "20%", "dir": RUN_V5 / "german_20pct"},
    {"dataset": "titanic", "noise_pct": "20%", "dir": RUN_V5 / "titanic_20pct"},
]
NOISE_TYPES = ["outlier", "rnd_label", "nnar", "mnar"]
PIPELINES   = ["p1a", "p2b"]
METHODS     = ["datascope", "cleanlab"]

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x2e, 0x4a)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1f, 0x77, 0xb4)   # DataScope blue
LIGHT  = RGBColor(0xf4, 0xf6, 0xf9)
GRAY   = RGBColor(0x55, 0x55, 0x55)

# Slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs: Presentation):
    """Return a truly blank slide."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _fill(shape, r, g, b):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def _rect(slide, left, top, w, h, r, g, b):
    shape = slide.shapes.add_shape(1, left, top, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    _fill(shape, r, g, b)
    shape.line.fill.background()
    return shape


def _textbox(slide, text, left, top, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def _add_image(slide, img_path: Path, left, top, w, h):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, w, h)


def _header_bar(slide, title: str, subtitle: str = ""):
    _rect(slide, 0, 0, W, Inches(1.1), 0x1a, 0x2e, 0x4a)
    _textbox(slide, title, Inches(0.35), Inches(0.12), Inches(10), Inches(0.55),
             size=24, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, subtitle, Inches(0.35), Inches(0.65), Inches(10), Inches(0.35),
                 size=13, color=RGBColor(0xbb, 0xcc, 0xdd))


def _load_curves(run_dir: Path, dataset: str, noise_type: str, pipeline: str):
    p = run_dir / "caches" / f"{dataset}__{noise_type}__{pipeline}" / "summary.json"
    if not p.exists():
        return None
    return json.loads(p.read_text())["curves"]


def _final(curves: dict, method: str):
    v = curves.get(method)
    if v is None:
        return None
    return float(v[-1]) if isinstance(v, list) else float(v)


def _best(curves: dict) -> tuple[str, float]:
    labels = {"datascope": "DataScope", "cleanlab": "CleanLab"}
    scores = {m: _final(curves, m) for m in METHODS if _final(curves, m) is not None}
    if not scores:
        return "—", 0.0
    b = max(scores, key=scores.get)
    return labels[b], scores[b]


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0x1a, 0x2e, 0x4a)
    # accent bar
    _rect(sl, 0, Inches(3.2), Inches(0.08), Inches(1.1), 0x1f, 0x77, 0xb4)
    _textbox(sl, "Label Cleaner",
             Inches(0.5), Inches(2.1), Inches(12), Inches(0.9),
             size=44, bold=True, color=WHITE)
    _textbox(sl, "Benchmarking Data-Cleaning Methods for Noisy Training Labels",
             Inches(0.5), Inches(3.1), Inches(11), Inches(0.6),
             size=20, color=RGBColor(0xbb, 0xcc, 0xdd))
    _textbox(sl, "Adult (20% noise)  ·  German Credit (20%)  ·  Titanic (20%)\n"
                 "DataScope · CleanLab · Random baseline  ·  4 noise types  ·  2 pipelines",
             Inches(0.5), Inches(3.85), Inches(11), Inches(0.7),
             size=14, color=RGBColor(0x88, 0xaa, 0xcc))
    _textbox(sl, "run_v5", Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.3),
             size=10, color=RGBColor(0x55, 0x77, 0x99), align=PP_ALIGN.RIGHT)


def slide_setup(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Experimental Setup")

    # Datasets column
    _textbox(sl, "Datasets", Inches(0.35), Inches(1.25), Inches(4), Inches(0.35),
             size=14, bold=True, color=NAVY)
    rows = [
        ("Adult",   "48 842 rows · 14 features · Sex",    "Income >50K",  "20%"),
        ("German",  "1 000 rows · 20 features · Sex",     "Credit risk",  "20%"),
        ("Titanic", "891 rows  · 7 features  · Sex",      "Survival",     "20%"),
    ]
    for i, (name, desc, task, noise) in enumerate(rows):
        y = Inches(1.65 + i * 0.9)
        _rect(sl, Inches(0.35), y, Inches(4.0), Inches(0.75), 0xFF, 0xFF, 0xFF)
        _textbox(sl, f"{name}  ({noise} noise)", Inches(0.5), y + Inches(0.06),
                 Inches(3.7), Inches(0.3), size=12, bold=True, color=NAVY)
        _textbox(sl, f"{desc}\nTask: {task}", Inches(0.5), y + Inches(0.33),
                 Inches(3.7), Inches(0.38), size=10, color=GRAY)

    # Noise types column
    _textbox(sl, "Noise Types", Inches(4.85), Inches(1.25), Inches(3.8), Inches(0.35),
             size=14, bold=True, color=NAVY)
    noise_rows = [
        ("Outlier",      "Feature values replaced with extreme value (100)\nCleaning: cap at 2σ"),
        ("Random Label", "Labels flipped uniformly at random\nCleaning: restore ground-truth"),
        ("NNAR",         "Labels flipped for protected subgroup only\nCleaning: restore ground-truth"),
        ("MNAR",         "Features set to NaN for protected subgroup\nCleaning: remove detected rows"),
    ]
    for i, (name, desc) in enumerate(noise_rows):
        y = Inches(1.65 + i * 0.9)
        _rect(sl, Inches(4.85), y, Inches(3.8), Inches(0.75), 0xFF, 0xFF, 0xFF)
        _textbox(sl, name, Inches(5.0), y + Inches(0.06), Inches(3.5), Inches(0.3),
                 size=12, bold=True, color=NAVY)
        _textbox(sl, desc, Inches(5.0), y + Inches(0.33), Inches(3.5), Inches(0.38),
                 size=9.5, color=GRAY)

    # Pipelines column
    _textbox(sl, "Pipelines", Inches(9.15), Inches(1.25), Inches(3.8), Inches(0.35),
             size=14, bold=True, color=NAVY)
    pipe_rows = [
        ("p1a", "KNN Imputer → StandardScaler → LogReg\nGeneral-purpose"),
        ("p2b", "PCA (n=10) → SelectKBest (k=8) → LogReg\nDimensionality reduction"),
    ]
    for i, (key, desc) in enumerate(pipe_rows):
        y = Inches(1.65 + i * 0.9)
        _rect(sl, Inches(9.15), y, Inches(3.8), Inches(0.75), 0xFF, 0xFF, 0xFF)
        _textbox(sl, key, Inches(9.3), y + Inches(0.06), Inches(3.5), Inches(0.3),
                 size=12, bold=True, color=NAVY)
        _textbox(sl, desc, Inches(9.3), y + Inches(0.33), Inches(3.5), Inches(0.38),
                 size=9.5, color=GRAY)



def slide_methods(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Cleaning Methods")

    methods = [
        ("DataScope",      "#1f77b4",
         "Ranks noisy samples by ascending Shapley importance — the samples that "
         "hurt test accuracy most are cleaned first. Uses NEIGHBOR approximation "
         "(KNN-based)."),
        ("CleanLab",       "#d62728",
         "OOF self_confidence scoring via cross_val_predict. "
         "Ranks ALL training samples; most suspicious (lowest confidence) first. "
         "Fully unsupervised — does not use ground-truth noisy positions."),
        ("DS Removal",     "#2ca02c",
         "DataScope ranking, but detected rows are removed from training instead "
         "of capped. Outlier experiments only — tests whether dropping beats "
         "repairing."),
        ("Random",         "#ff7f0e",
         "Shuffles noisy positions uniformly at random. Mean ± 1σ over 3 seeds. "
         "Baseline that isolates the ordering benefit of learned methods."),
    ]

    cols = 4
    box_w = Inches(3.1)
    box_h = Inches(1.55)
    for i, (name, hex_col, desc) in enumerate(methods):
        row, col = divmod(i, cols)
        left = Inches(0.2 + col * 3.3)
        top  = Inches(1.25 + row * 1.65)
        r, g, b = int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16)
        # accent strip
        _rect(sl, left, top, Inches(0.06), box_h, r, g, b)
        _rect(sl, left + Inches(0.06), top, box_w - Inches(0.06), box_h, 0xFF, 0xFF, 0xFF)
        _textbox(sl, name, left + Inches(0.15), top + Inches(0.08),
                 box_w - Inches(0.2), Inches(0.3), size=12, bold=True, color=NAVY)
        _textbox(sl, desc, left + Inches(0.15), top + Inches(0.38),
                 box_w - Inches(0.2), Inches(1.1), size=8.5, color=GRAY)


def _step_box(slide, number: str, text: str, left, top, w=Inches(5.8), h=Inches(0.62),
              accent_rgb=(0x1f, 0x77, 0xb4)):
    """Numbered step box: accent circle + text."""
    _rect(slide, left, top, Inches(0.48), h, *accent_rgb)
    _textbox(slide, number, left, top + Inches(0.1), Inches(0.48), h - Inches(0.2),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, left + Inches(0.48), top, w - Inches(0.48), h, 0xFF, 0xFF, 0xFF)
    _textbox(slide, text, left + Inches(0.58), top + Inches(0.08),
             w - Inches(0.65), h - Inches(0.16), size=10.5, color=NAVY)


def _pro_con(slide, pros: list[str], cons: list[str], left_x, top_y):
    """Two-column strengths / weaknesses block."""
    col_w = Inches(3.2)
    # Strengths
    _rect(slide, left_x, top_y, col_w, Inches(0.34), 0x2c, 0xa0, 0x2c)
    _textbox(slide, "Strengths", left_x + Inches(0.08), top_y + Inches(0.05),
             col_w, Inches(0.26), size=11, bold=True, color=WHITE)
    for i, p in enumerate(pros):
        _textbox(slide, f"+ {p}", left_x + Inches(0.08),
                 top_y + Inches(0.36 + i * 0.44), col_w, Inches(0.42), size=9.5, color=NAVY)
    # Weaknesses
    rx = left_x + col_w + Inches(0.15)
    _rect(slide, rx, top_y, col_w, Inches(0.34), 0xd6, 0x27, 0x28)
    _textbox(slide, "Weaknesses", rx + Inches(0.08), top_y + Inches(0.05),
             col_w, Inches(0.26), size=11, bold=True, color=WHITE)
    for i, c in enumerate(cons):
        _textbox(slide, f"− {c}", rx + Inches(0.08),
                 top_y + Inches(0.36 + i * 0.44), col_w, Inches(0.42), size=9.5, color=NAVY)


def slide_method_datascope(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Method: DataScope", "Shapley-value importance ranking")

    # Left: how it works
    _textbox(sl, "How it works", Inches(0.35), Inches(1.2), Inches(6.5), Inches(0.35),
             size=13, bold=True, color=NAVY)
    _textbox(sl,
             "Assigns each training sample a Shapley value measuring its marginal "
             "contribution to test accuracy. Samples with the most negative contribution "
             "— those that actively hurt accuracy — are ranked first and cleaned.\n\n"
             "Uses the NEIGHBOR approximation (KNN-based) to avoid enumerating all "
             "subsets, making it tractable up to ~50k rows.",
             Inches(0.35), Inches(1.58), Inches(6.5), Inches(1.6), size=11, color=GRAY)

    _textbox(sl, "Algorithm", Inches(0.35), Inches(3.25), Inches(6.5), Inches(0.35),
             size=13, bold=True, color=NAVY)
    steps = [
        ("1", "Fit the full pipeline on the noisy training set"),
        ("2", "Compute Shapley importances via ShapleyImportance(NEIGHBOR).score(X_test)"),
        ("3", "Sort noisy positions by ascending importance (most harmful first)"),
        ("4", "Incrementally apply action_fn to top-k% · measure accuracy at each step"),
    ]
    for i, (n, t) in enumerate(steps):
        _step_box(sl, n, t, Inches(0.35), Inches(3.65 + i * 0.68))

    _pro_con(sl,
             ["Theoretically grounded — only method with axiomatic fairness guarantees",
              "Consistent across all noise types; never catastrophically fails",
              "Works on both label and feature noise"],
             ["Ranking computed once — becomes stale as noise is cleaned",
              "No label uncertainty signal — treats mislabelled = unusual samples equally",
              "Monte Carlo variant: 56+ min on adult (50 iterations used as default)"],
             Inches(7.1), Inches(1.2))

    # Formula box
    _rect(sl, Inches(7.1), Inches(4.55), Inches(6.0), Inches(1.35), 0xFF, 0xFF, 0xFF)
    _textbox(sl, "Shapley value (exact):",
             Inches(7.25), Inches(4.62), Inches(5.7), Inches(0.3), size=10, bold=True, color=NAVY)
    _textbox(sl,
             "φᵢ = Σ  |S|!(n−|S|−1)!/n!  ·  [v(S∪{i}) − v(S)]\n"
             "     S⊆N\\{i}\n\n"
             "v(S) = accuracy of model trained on subset S",
             Inches(7.25), Inches(4.96), Inches(5.7), Inches(0.85),
             size=10, color=GRAY)


def slide_method_cleanlab(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Method: CleanLab", "Self-confidence ranking")

    _textbox(sl, "How it works", Inches(0.35), Inches(1.2), Inches(6.5), Inches(0.35),
             size=13, bold=True, color=NAVY)
    _textbox(sl,
             "Out-of-fold (OOF) predicted probabilities via 5-fold cross-validation. "
             "Self-confidence = P(current label correct | x). Lower = more suspicious. "
             "Fully unsupervised — does not need ground-truth noisy positions.",
             Inches(0.35), Inches(1.58), Inches(6.5), Inches(1.2), size=11, color=GRAY)

    _textbox(sl, "Algorithm", Inches(0.35), Inches(3.0), Inches(6.5), Inches(0.35),
             size=13, bold=True, color=NAVY)
    cl_steps = [
        ("1", "cross_val_predict(cv=5, method='predict_proba') on full training set"),
        ("2", "find_label_issues(ranked_by='self_confidence') → ranked suspicious list"),
        ("3", "Apply action_fn to top-k% most suspicious samples"),
        ("4", "Measure accuracy on the clean held-out test set at each proportion"),
    ]
    for i, (n, t) in enumerate(cl_steps):
        _step_box(sl, n, t, Inches(0.35), Inches(3.4 + i * 0.68),
                  accent_rgb=(0xd6, 0x27, 0x28))

    _pro_con(sl,
             ["Unsupervised — no ground-truth noisy positions needed",
              "Strong on random label noise — self_confidence ranks flips reliably",
              "Cheap: one 5-fold CV pass, no importance computation"],
             ["Degrades at high noise rates — OOF probs become noisy themselves",
              "Less effective on feature noise (label correct, feature corrupted)",
              "Ranks all samples — top-k includes false positives on clean data"],
             Inches(7.1), Inches(1.2))


def slide_heatmap(prs: Presentation):
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Cross-Dataset Comparison",
                "Mean accuracy improvement over baseline (averaged across both pipelines) · Green = better")
    img = COMBINED / "comparison_heatmap.png"
    _add_image(sl, img, Inches(1.5), Inches(1.2), Inches(10.3), Inches(6.0))


def slide_dataset_grid(prs: Presentation, run: dict):
    ds     = run["dataset"]
    noise  = run["noise_pct"]
    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, f"{ds.capitalize()} — {noise} noise · All noise types",
                "Rows: p1a / p2b  ·  Cols: outlier / rnd_label / nnar / mnar")
    img = run["dir"] / "figures" / f"{ds}__all_noise_types.png"
    _add_image(sl, img, Inches(0.15), Inches(1.15), Inches(13.0), Inches(6.2))


def slide_findings(prs: Presentation):
    # Collect best numbers from run_v5 caches
    findings = []
    for run in RUNS:
        ds = run["dataset"]
        for nt in NOISE_TYPES:
            for pl in PIPELINES:
                c = _load_curves(run["dir"], ds, nt, pl)
                if c is None:
                    continue
                bl  = c.get("baseline", 0)
                name, val = _best(c)
                if val - bl > 0.01:
                    findings.append((ds, nt, pl, name, bl, val, val - bl))

    # Sort by improvement descending, take top 8
    findings.sort(key=lambda x: x[6], reverse=True)
    top = findings[:8]

    sl = _blank(prs)
    _rect(sl, 0, 0, W, H, 0xf4, 0xf6, 0xf9)
    _header_bar(sl, "Key Results — Top Accuracy Gains",
                "Configurations where cleaning improved accuracy most over baseline")

    col_headers = ["Dataset", "Noise", "Pipeline", "Best Method", "Baseline", "Final", "Δ"]
    col_widths  = [Inches(1.4), Inches(1.4), Inches(1.1), Inches(2.2), Inches(1.2), Inches(1.2), Inches(1.0)]
    col_starts  = [Inches(0.3)]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    row_h  = Inches(0.48)
    top_y  = Inches(1.2)

    # Header row
    _rect(sl, Inches(0.3), top_y, sum(col_widths), row_h, 0x1a, 0x2e, 0x4a)
    for j, (hdr, x) in enumerate(zip(col_headers, col_starts)):
        _textbox(sl, hdr, x + Inches(0.05), top_y + Inches(0.1),
                 col_widths[j], row_h, size=11, bold=True, color=WHITE)

    for i, (ds, nt, pl, name, bl, val, delta) in enumerate(top):
        y   = top_y + row_h + i * row_h
        bg  = (0xFF, 0xFF, 0xFF) if i % 2 == 0 else (0xeb, 0xf0, 0xf7)
        _rect(sl, Inches(0.3), y, sum(col_widths), row_h, *bg)
        cells = [ds, nt, pl, name, f"{bl:.4f}", f"{val:.4f}", f"+{delta:.4f}"]
        for j, (cell, x) in enumerate(zip(cells, col_starts)):
            color = RGBColor(0x1f, 0x77, 0xb4) if j == 6 else NAVY
            bold  = j == 6
            _textbox(sl, cell, x + Inches(0.05), y + Inches(0.1),
                     col_widths[j], row_h, size=10, bold=bold, color=color)



# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = _prs()

    slide_title(prs)
    slide_setup(prs)
    slide_methods(prs)
    slide_method_datascope(prs)
    slide_method_cleanlab(prs)
    slide_heatmap(prs)
    for run in RUNS:
        slide_dataset_grid(prs, run)
    slide_findings(prs)

    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    raise SystemExit(main())
